"""
文档处理服务模块
"""
import os
import io
import hashlib
import mimetypes
from pathlib import Path
from typing import Optional, List, Dict, Any, BinaryIO
import PyPDF2
import docx
import pandas as pd
from pptx import Presentation
import fitz  # PyMuPDF
from PIL import Image
import markdown
from bs4 import BeautifulSoup
# import textract  # 暂时注释掉，避免安装问题
# from sentence_transformers import SentenceTransformer  # 暂时注释掉，避免大型依赖
# import numpy as np
from app.core.config import settings
from app.core.logging import logger


class DocumentProcessor:
    """文档处理器 - 支持多种格式的文档解析和向量化"""

    def __init__(self):
        # 初始化向量化模型 - 暂时禁用，避免大型依赖
        # self.embedding_model = SentenceTransformer('shibing624/text2vec-base-chinese')
        self.embedding_model = None  # 暂时设为None
        self.supported_formats = {
            '.pdf', '.docx', '.doc', '.txt', '.md', '.html', '.htm',
            '.xlsx', '.xls', '.csv', '.pptx', '.ppt', '.jpg', '.jpeg',
            '.png', '.gif', '.bmp', '.tiff'
        }

    @staticmethod
    def _to_list(vector):
        """安全地将向量转换为list格式

        Args:
            vector: numpy数组、list或其他可迭代对象

        Returns:
            list或None
        """
        if vector is None:
            return None
        if isinstance(vector, list):
            return vector
        elif hasattr(vector, 'tolist'):
            # numpy数组或torch张量
            return vector.tolist()
        else:
            # 其他可迭代对象
            try:
                return list(vector)
            except:
                return None

    async def process_document(
        self,
        file_path: str,
        filename: str,
        chunk_size: int = 512,
        chunk_overlap: int = 50
    ) -> Dict[str, Any]:
        """
        处理文档并生成向量

        Args:
            file_path: 文件路径
            filename: 文件名
            chunk_size: 分块大小
            chunk_overlap: 分块重叠

        Returns:
            处理结果字典
        """
        try:
            # 1. 检测文件类型
            file_ext = Path(filename).suffix.lower()
            if file_ext not in self.supported_formats:
                raise ValueError(f"不支持的文件格式: {file_ext}")

            # 2. 计算文件哈希
            file_hash = await self._calculate_file_hash(file_path)

            # 3. 提取文本内容
            text_content = await self._extract_text(file_path, file_ext)

            # 4. 提取元数据
            metadata = await self._extract_metadata(file_path, file_ext)

            # 5. 文本分块
            chunks = await self._chunk_text(
                text_content,
                chunk_size=chunk_size,
                chunk_overlap=chunk_overlap
            )

            # 6. 生成向量
            embeddings = await self._generate_embeddings(chunks)

            # 7. 处理结果
            result = {
                'file_hash': file_hash,
                'file_type': file_ext,
                'total_length': len(text_content),
                'chunk_count': len(chunks),
                'metadata': metadata,
                'chunks': chunks,
                'embeddings': embeddings,  # embeddings已经是List[List[float]]格式，不需要转换
                'processing_status': 'success'
            }

            logger.info(f"文档处理完成: {filename}, 分块数: {len(chunks)}")
            return result

        except Exception as e:
            logger.error(f"文档处理失败 {filename}: {str(e)}")
            return {
                'processing_status': 'error',
                'error_message': str(e),
                'file_hash': None,
                'chunks': [],
                'embeddings': None
            }

    async def _calculate_file_hash(self, file_path: str) -> str:
        """计算文件哈希值"""
        hash_md5 = hashlib.md5()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_md5.update(chunk)
        return hash_md5.hexdigest()

    async def _extract_text(self, file_path: str, file_ext: str) -> str:
        """根据文件类型提取文本"""
        try:
            if file_ext == '.pdf':
                return await self._extract_pdf_text(file_path)
            elif file_ext in ['.docx', '.doc']:
                return await self._extract_docx_text(file_path)
            elif file_ext in ['.txt', '.md']:
                return await self._extract_text_file(file_path)
            elif file_ext in ['.html', '.htm']:
                return await self._extract_html_text(file_path)
            elif file_ext in ['.xlsx', '.xls']:
                return await self._extract_excel_text(file_path)
            elif file_ext == '.csv':
                return await self._extract_csv_text(file_path)
            elif file_ext == '.pptx':
                return await self._extract_pptx_text(file_path)
            elif file_ext == '.ppt':
                # 旧版PPT格式不支持，给出明确提示
                raise ValueError(
                    "不支持旧版.ppt格式。请使用PowerPoint将文件另存为.pptx格式后重新上传。"
                )
            elif file_ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff']:
                return await self._extract_image_text(file_path)
            else:
                raise ValueError(f"不支持的文件格式: {file_ext}")
        except Exception as e:
            logger.error(f"文本提取失败 {file_path}: {str(e)}")
            raise

    async def _validate_pdf_file(self, file_path: str) -> tuple[bool, str]:
        """验证PDF文件是否可读

        Returns:
            (is_valid, error_message): 是否有效和错误消息
        """
        try:
            # 使用PyMuPDF验证
            doc = fitz.open(file_path)
            page_count = doc.page_count
            is_encrypted = doc.is_encrypted
            doc.close()

            if page_count == 0:
                return False, "PDF文件没有页面"
            if is_encrypted:
                return False, "PDF文件已加密，无法处理"

            logger.info(f"PDF验证通过: {page_count}页")
            return True, ""

        except Exception as e:
            error_msg = str(e)
            if "format" in error_msg.lower():
                return False, f"PDF格式错误或文件损坏: {error_msg}"
            elif "password" in error_msg.lower():
                return False, "PDF文件需要密码"
            else:
                return False, f"PDF文件无法打开: {error_msg}"

    async def _extract_pdf_text(self, file_path: str) -> str:
        """提取PDF文本 - 使用PyMuPDF"""
        # 先验证PDF文件
        is_valid, error_msg = await self._validate_pdf_file(file_path)
        if not is_valid:
            raise ValueError(f"PDF验证失败: {error_msg}")

        try:
            doc = fitz.open(file_path)
            text = ""
            for page_num in range(doc.page_count):
                page = doc[page_num]
                page_text = page.get_text()
                text += page_text

            doc.close()

            if not text.strip():
                logger.warning(f"PDF文件 {file_path} 提取的文本为空，可能是扫描版PDF")

            return text.strip()

        except Exception as e:
            logger.error(f"PyMuPDF文本提取失败: {str(e)}", exc_info=True)
            # 备用方案：使用PyPDF2
            try:
                logger.info("尝试使用PyPDF2作为备用方案")
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    text = ""
                    for page in pdf_reader.pages:
                        text += page.extract_text()

                if not text.strip():
                    raise ValueError("PDF文本提取为空，可能是扫描版PDF或图片格式")

                return text.strip()
            except Exception as backup_error:
                logger.error(f"PyPDF2备用方案也失败: {str(backup_error)}", exc_info=True)
                raise ValueError(f"无法解析PDF文件: PyMuPDF失败({str(e)}), PyPDF2失败({str(backup_error)})")

    async def _extract_docx_text(self, file_path: str) -> str:
        """提取Word文档文本"""
        try:
            doc = docx.Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text.strip()
        except Exception as e:
            logger.error(f"Word文档文本提取失败: {str(e)}")
            raise

    async def _extract_text_file(self, file_path: str) -> str:
        """提取纯文本文件"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read().strip()
        except UnicodeDecodeError:
            # 尝试其他编码
            try:
                with open(file_path, 'r', encoding='gbk') as file:
                    return file.read().strip()
            except:
                with open(file_path, 'r', encoding='latin-1') as file:
                    return file.read().strip()

    async def _extract_html_text(self, file_path: str) -> str:
        """提取HTML文本"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                html_content = file.read()
            soup = BeautifulSoup(html_content, 'html.parser')
            return soup.get_text(separator='\n', strip=True)
        except Exception as e:
            logger.error(f"HTML文本提取失败: {str(e)}")
            raise

    async def _extract_excel_text(self, file_path: str) -> str:
        """提取Excel文本"""
        try:
            df = pd.read_excel(file_path)
            text = ""
            for index, row in df.iterrows():
                row_text = "\t".join([str(cell) if pd.notna(cell) else "" for cell in row])
                text += row_text + "\n"
            return text.strip()
        except Exception as e:
            logger.error(f"Excel文本提取失败: {str(e)}")
            raise

    async def _extract_csv_text(self, file_path: str) -> str:
        """提取CSV文本"""
        try:
            df = pd.read_csv(file_path)
            text = df.to_string(index=False)
            return text.strip()
        except Exception as e:
            logger.error(f"CSV文本提取失败: {str(e)}")
            raise

    async def _extract_pptx_text(self, file_path: str) -> str:
        """提取PowerPoint文本"""
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text.strip()
        except Exception as e:
            logger.error(f"PowerPoint文本提取失败: {str(e)}")
            raise

    async def _extract_image_text(self, file_path: str) -> str:
        """提取图片文本（OCR）- 简化版本"""
        try:
            # 这里可以集成OCR服务如Tesseract或云OCR API
            # 目前返回文件名作为占位符
            filename = Path(file_path).name
            return f"[图片文件: {filename} - 需要OCR识别]"
        except Exception as e:
            logger.error(f"图片文本提取失败: {str(e)}")
            raise

    async def _extract_metadata(self, file_path: str, file_ext: str) -> Dict[str, Any]:
        """提取文档元数据"""
        try:
            file_stat = os.stat(file_path)
            metadata = {
                'filename': Path(file_path).name,
                'file_size': file_stat.st_size,
                'created_time': file_stat.st_ctime,
                'modified_time': file_stat.st_mtime,
                'mime_type': mimetypes.guess_type(file_path)[0],
                'file_extension': file_ext
            }

            # 特定格式的额外元数据
            if file_ext == '.pdf':
                try:
                    doc = fitz.open(file_path)
                    metadata.update({
                        'page_count': doc.page_count,
                        'pdf_version': doc.pdf_version(),
                        'title': doc.metadata.get('title', ''),
                        'author': doc.metadata.get('author', ''),
                        'subject': doc.metadata.get('subject', ''),
                        'creator': doc.metadata.get('creator', ''),
                        'producer': doc.metadata.get('producer', ''),
                        'creation_date': doc.metadata.get('creationDate', ''),
                        'modification_date': doc.metadata.get('modDate', '')
                    })
                    doc.close()
                except:
                    pass
            elif file_ext == '.docx':
                try:
                    doc = docx.Document(file_path)
                    metadata.update({
                        'paragraph_count': len(doc.paragraphs),
                        'word_count': len(doc.paragraphs) * 10,  # 估算
                        'author': doc.core_properties.author or '',
                        'title': doc.core_properties.title or '',
                        'subject': doc.core_properties.subject or '',
                        'created': doc.core_properties.created,
                        'modified': doc.core_properties.modified
                    })
                except:
                    pass
            elif file_ext in ['.xlsx', '.xls']:
                try:
                    df = pd.read_excel(file_path)
                    metadata.update({
                        'row_count': len(df),
                        'column_count': len(df.columns),
                        'sheet_names': ['Sheet1']  # 简化版本
                    })
                except:
                    pass
            elif file_ext == '.csv':
                try:
                    df = pd.read_csv(file_path)
                    metadata.update({
                        'row_count': len(df),
                        'column_count': len(df.columns),
                        'has_header': True  # 默认假设有标题
                    })
                except:
                    pass

            return metadata

        except Exception as e:
            logger.error(f"元数据提取失败: {str(e)}")
            return {
                'filename': Path(file_path).name,
                'file_extension': file_ext,
                'error': str(e)
            }

    async def _chunk_text(
        self,
        text: str,
        chunk_size: int = 512,
        chunk_overlap: int = 50
    ) -> List[Dict[str, Any]]:
        """文本分块"""
        if not text.strip():
            return []

        chunks = []
        text_length = len(text)

        # 按段落分割
        paragraphs = text.split('\n')
        current_chunk = ""
        chunk_index = 0

        for para in paragraphs:
            para = para.strip()
            if not para:
                continue

            # 如果当前块加上新段落不超过大小限制
            if len(current_chunk) + len(para) + 1 <= chunk_size:
                if current_chunk:
                    current_chunk += "\n" + para
                else:
                    current_chunk = para
            else:
                # 保存当前块
                if current_chunk:
                    chunks.append({
                        'chunk_index': chunk_index,
                        'content': current_chunk.strip(),
                        'start_position': sum(len(c['content']) for c in chunks),
                        'end_position': sum(len(c['content']) for c in chunks) + len(current_chunk),
                        'word_count': len(current_chunk.split()),
                        'char_count': len(current_chunk)
                    })
                    chunk_index += 1

                # 开始新块
                current_chunk = para

                # 如果单个段落太长，需要进一步分割
                if len(para) > chunk_size:
                    words = para.split()
                    temp_chunk = ""
                    for word in words:
                        if len(temp_chunk) + len(word) + 1 <= chunk_size:
                            if temp_chunk:
                                temp_chunk += " " + word
                            else:
                                temp_chunk = word
                        else:
                            if temp_chunk:
                                chunks.append({
                                    'chunk_index': chunk_index,
                                    'content': temp_chunk.strip(),
                                    'start_position': sum(len(c['content']) for c in chunks),
                                    'end_position': sum(len(c['content']) for c in chunks) + len(temp_chunk),
                                    'word_count': len(temp_chunk.split()),
                                    'char_count': len(temp_chunk)
                                })
                                chunk_index += 1
                            temp_chunk = word
                    current_chunk = temp_chunk

        # 处理最后一个块
        if current_chunk:
            chunks.append({
                'chunk_index': chunk_index,
                'content': current_chunk.strip(),
                'start_position': sum(len(c['content']) for c in chunks),
                'end_position': sum(len(c['content']) for c in chunks) + len(current_chunk),
                'word_count': len(current_chunk.split()),
                'char_count': len(current_chunk)
            })

        # 添加重叠信息
        for i, chunk in enumerate(chunks):
            if i > 0:
                chunk['previous_chunk_overlap'] = chunks[i-1]['content'][-chunk_overlap:] if len(chunks[i-1]['content']) > chunk_overlap else chunks[i-1]['content']
            if i < len(chunks) - 1:
                chunk['next_chunk_overlap'] = chunks[i+1]['content'][:chunk_overlap] if len(chunks[i+1]['content']) > chunk_overlap else chunks[i+1]['content']

        return chunks

    async def _generate_embeddings(self, chunks: List[Dict[str, Any]]) -> Optional[List[List[float]]]:
        """生成文本向量 - 从数据库读取配置并使用对应的Embedding API"""
        try:
            if not chunks:
                return None

            # 从数据库读取embedding配置
            from sqlalchemy import select
            from app.models.system_config import SystemConfig
            from app.db.session import AsyncSessionLocal
            
            async with AsyncSessionLocal() as session:
                result = await session.execute(
                    select(SystemConfig).where(
                        SystemConfig.category == "embedding",
                        SystemConfig.is_active == True
                    )
                )
                configs = result.scalars().all()
                
                # 解析配置
                embedding_config = {}
                for config in configs:
                    key = config.config_key.replace("embedding_", "")
                    embedding_config[key] = config.config_value
            
            provider = embedding_config.get('provider', 'dashscope')
            model = embedding_config.get('model', 'text-embedding-v2')
            api_key = embedding_config.get('api_key', '')
            
            if not api_key:
                logger.error(f"Embedding API密钥未配置（供应商：{provider}）")
                return None
            
            logger.info(f"使用Embedding配置 - 供应商：{provider}, 模型：{model}")
            
            # 根据供应商选择对应的API
            if provider in ['dashscope', 'qwen']:
                # 使用阿里DashScope
                import dashscope
                from dashscope import TextEmbedding
                
                dashscope.api_key = api_key
                embeddings = []
                
                for chunk in chunks:
                    content = chunk.get('content', '')
                    if not content or len(content.strip()) == 0:
                        embeddings.append([0.0] * 1536)
                        continue
                    
                    try:
                        # 根据模型名称选择对应的模型
                        model_enum = TextEmbedding.Models.text_embedding_v2
                        if 'v1' in model.lower():
                            model_enum = TextEmbedding.Models.text_embedding_v1
                        elif 'v3' in model.lower():
                            model_enum = TextEmbedding.Models.text_embedding_v3
                        
                        response = TextEmbedding.call(
                            model=model_enum,
                            input=content[:2000]
                        )
                        
                        if response.status_code == 200:
                            embedding = response.output['embeddings'][0]['embedding']
                            embeddings.append(embedding)
                        else:
                            logger.error(f"DashScope API调用失败: {response.message}")
                            embeddings.append([0.0] * 1536)
                            
                    except Exception as e:
                        logger.error(f"单个chunk向量生成失败: {str(e)}")
                        embeddings.append([0.0] * 1536)
                
                logger.info(f"成功生成 {len(embeddings)} 个向量（使用{provider}/{model}）")
                return embeddings
            
            else:
                logger.error(f"不支持的Embedding供应商: {provider}")
                return None

        except Exception as e:
            logger.error(f"向量生成失败: {str(e)}")
            import traceback
            logger.error(traceback.format_exc())
            return None

    async def get_supported_formats(self) -> List[str]:
        """获取支持的文件格式列表"""
        return list(self.supported_formats)

    async def validate_file(self, file_path: str, max_size_mb: int = 100) -> Dict[str, Any]:
        """验证文件是否符合处理要求"""
        try:
            if not os.path.exists(file_path):
                return {'valid': False, 'error': '文件不存在'}

            file_size = os.path.getsize(file_path)
            max_size_bytes = max_size_mb * 1024 * 1024

            if file_size > max_size_bytes:
                return {
                    'valid': False,
                    'error': f'文件大小超过限制 ({max_size_mb}MB)'
                }

            file_ext = Path(file_path).suffix.lower()
            if file_ext not in self.supported_formats:
                return {
                    'valid': False,
                    'error': f'不支持的文件格式: {file_ext}'
                }

            return {
                'valid': True,
                'file_size': file_size,
                'file_type': file_ext,
                'estimated_processing_time': file_size / (1024 * 1024) * 2  # 估算处理时间
            }

        except Exception as e:
            return {'valid': False, 'error': str(e)}


# 全局文档处理器实例
document_processor = DocumentProcessor()